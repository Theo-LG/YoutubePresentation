from pptx import Presentation
from pptx.util import Inches
import logging


def save_presentation(slides_number, subtitles_list):
    """Given a list of frames indices, create a presentation

    Args:
        slides_number (int): Number of frames saved
        subtitles_list (list): List of text to save as notes for each slides
    """
    logging.info("Creating powerpoint...")
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_slide_layout = prs.slide_layouts[6]

    for i in range(0, slides_number):
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(
            "Frames/image" + str(i) + ".jpg",
            0,
            0,
            height=prs.slide_height,
        )
        # Add subtitles as note
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = subtitles_list[i]
    prs.save("YouTubeSlides.pptx")
    logging.info("PowerPoint saved")
